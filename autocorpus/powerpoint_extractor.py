"""This module provides functionality to extract text from PowerPoint presentations.

Functions:
- extract_text_from_shape: Extracts text from individual shapes in a slide.
- get_powerpoint_text: Extracts all text from a PowerPoint presentation file.
"""

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape


def extract_text_from_shape(shape: BaseShape) -> list[str]:
    """Extract text from a PowerPoint shape.

    Parameters:
        shape : BaseShape
            The shape object from which to extract text.

    Returns:
        list[str]
            A list of strings containing the text extracted from the shape.
    """
    text = []

    # If the shape is a group, recurse into each sub-shape
    if isinstance(shape, GroupShape):
        for sub_shape in shape.shapes:
            text.extend(extract_text_from_shape(sub_shape))

    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:  # type: ignore
        for paragraph in shape.text_frame.paragraphs:  # type: ignore
            for run in paragraph.runs:
                if run.text.strip():
                    text.append(run.text)

    return text


def get_powerpoint_text(filename: str) -> list[str]:
    """Extract all text from a PowerPoint presentation.

    Parameters:
        filename : str
            The path to the PowerPoint (.pptx) file.

    Returns:
        list[str]
            A list of strings containing all the text found in the presentation.
    """
    text = []
    pres = Presentation(filename)
    for slide in pres.slides:
        for shape in slide.shapes:
            text.extend(extract_text_from_shape(shape))
    return text
